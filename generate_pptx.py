#!/usr/bin/env python3
"""
Generate a professional PPTX presentation from slides and speaker notes.
Embeds output images and speaker notes into the slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define slide data (title, bullet points, image_path, speaker_notes)
slides_data = [
    {
        "title": "Professional Unemployment Analysis",
        "subtitle": "COVID-19 Impact & Forecasting (India 2019–2020)",
        "bullets": [
            "Objective: Quantify lockdown impact on employment",
            "Identify seasonal patterns and temporal trends",
            "Produce short-term forecasts for policy guidance",
        ],
        "image": None,
        "speaker_notes": "Hello — I'm presenting a professional analysis of unemployment in India for 2019–2020. The goals: quantify the lockdown impact, identify seasonal/temporal patterns, and produce short-term forecasts to inform policy.",
    },
    {
        "title": "Data & Cleaning",
        "subtitle": "Multi-source Integration",
        "bullets": [
            "Source files: Unemployment in India.xls, Unemployment_Rate_upto_11_2020.xls",
            "Cleaning: header normalization, date parsing (day-first), numeric coercion",
            "Auto-discovery and deduplication of unemployment records",
            "Final dataset: 3,494 unique observations (May 2019 – Oct 2020)",
        ],
        "image": None,
        "speaker_notes": "We used two main source files. Many files have inconsistent headers; we normalize columns, parse dates with day-first, coerce numeric types, and drop invalid rows. The pipeline auto-discovers files with 'unemployment' in their names and concatenates them, deduplicating records.",
    },
    {
        "title": "Methodology",
        "subtitle": "Statistical & Time-Series Approach",
        "bullets": [
            "Feature engineering: month, year, COVID period (Pre-, Lockdown, Post-)",
            "Aggregation levels: national, state/region, rural/urban",
            "Statistical test: Welch's t-test (unequal variances)",
            "Forecasting: SARIMA with ARIMA & AR(1) OLS fallbacks for robustness",
        ],
        "image": None,
        "speaker_notes": "We created time features and defined COVID periods: Pre-COVID (before 2020-03-01), Lockdown (Mar–Jun 2020), Post-Lockdown (after Jul 2020). For inference we use Welch's t-test to compare means, and for forecasting we attempt SARIMA with fallbacks to ARIMA and a simple AR(1) OLS when models don't converge.",
    },
    {
        "title": "Key Visuals",
        "subtitle": "Observed Trends & Patterns",
        "bullets": [
            "National unemployment trend: clear spike during lockdown",
            "Rural vs Urban: Urban more volatile; Rural more stable",
            "Seasonality: recurring monthly patterns (informal sector sensitivity)",
        ],
        "image": "/home/abidoye/Desktop/CodeAlpha/national_trend.png",
        "speaker_notes": "This shows a spike in unemployment during the lockdown period. Rural and urban dynamics differ; urban is more volatile. There are recurring monthly patterns indicating informal-sector seasonality.",
    },
    {
        "title": "Statistical Results & COVID Impact",
        "subtitle": "Quantified Lockdown Effect",
        "bullets": [
            "Welch t-test: t = -16.64, p < 0.001 (highly significant)",
            "Pre-COVID avg: ~9.2% | Lockdown avg: ~16.8%",
            "COVID-19 increased unemployment by ~81% overall",
            "Rural +99.9% | Urban +77.8% (differential impact)",
        ],
        "image": None,
        "speaker_notes": "We aggregated average unemployment across defined COVID periods. The Welch t-test rejects the null (p < 0.05), indicating a statistically significant increase during lockdown. The pivot table quantifies absolute and percentage changes by area.",
    },
    {
        "title": "Forecasting & Model Selection",
        "subtitle": "SARIMA with Robust Fallbacks",
        "bullets": [
            "Primary: stabilized SARIMA(0,1,1)×(0,1,1,12) for national series",
            "Fallback 1: non-seasonal ARIMA(0,1,1) if SARIMA fails",
            "Fallback 2: AR(1) via OLS if both fail (always converges)",
            "6-month forecast: unemployment returning to 7–8% by early 2021",
        ],
        "image": "/home/abidoye/Desktop/CodeAlpha/sarima_forecast.png",
        "speaker_notes": "We fit a stabilized SARIMA model where possible. If the SARIMA fit was unstable, we tried ARIMA, then AR(1) OLS. The final SARIMA forecast is shown; interpret with caution due to sample period and possible structural breaks from the pandemic.",
    },
    {
        "title": "Policy Recommendations & Limitations",
        "subtitle": "Insights & Next Steps",
        "bullets": [
            "Recommendations: counter-cyclical urban programs, rural safety nets, digital skills",
            "Limitations: sample length, informal sector measurement, aggregation choices",
            "Future work: add covariates, state-level hierarchical modeling, real-time updates",
        ],
        "image": None,
        "speaker_notes": "Recommend counter-cyclical urban employment programs, strengthen rural safety nets, and invest in remote-work skills. Limitations: data representativeness, sample length, and model assumptions. Next steps: add covariates and state-level hierarchical modeling.",
    },
]

# Helper to add text box
def add_text_frame(shape, text, font_size=18, bold=False, color=RGBColor(0, 0, 0)):
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color

# Create slides
for i, slide_info in enumerate(slides_data):
    # Use blank layout and add shapes manually for flexibility
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    add_text_frame(title_box, slide_info["title"], font_size=44, bold=True, color=RGBColor(0, 51, 102))

    # Add subtitle if present
    if "subtitle" in slide_info:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
        add_text_frame(subtitle_box, slide_info["subtitle"], font_size=20, color=RGBColor(100, 100, 100))

    # Add bullets
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(4.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    for j, bullet in enumerate(slide_info["bullets"]):
        if j == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    # Add image if provided
    if slide_info["image"] and Path(slide_info["image"]).exists():
        try:
            # Determine position: right side if there are bullets, full slide if no bullets
            img_path = slide_info["image"]
            slide.shapes.add_picture(img_path, Inches(4.5), Inches(1.6), width=Inches(5))
        except Exception as e:
            print(f"Warning: Could not add image {slide_info['image']}: {e}")

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = slide_info["speaker_notes"]

print("✓ Presentation generated: unemployment_analysis.pptx")
prs.save("/home/abidoye/Desktop/CodeAlpha/unemployment_analysis.pptx")
